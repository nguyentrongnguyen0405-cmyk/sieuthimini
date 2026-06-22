import collections
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

# Khởi tạo presentation
prs = Presentation()

slides_data = [
    {
        "title": "Báo Cáo Bài Tập Lớn Công Nghệ Phần Mềm",
        "content": "Phân tích & Thiết kế Hệ thống POS & Cửa hàng Trực tuyến MiniMart POS\n\nLớp học: 23DTHC5 - 23DTHC4\nNhóm thực hiện: Nhóm [Tên nhóm]\nDanh sách thành viên & Vai trò:\n- [Họ tên 1] - Trưởng nhóm\n- [Họ tên 2] - Thành viên\n- [Họ tên 3] - Thành viên"
    },
    {
        "title": "Thực Trạng & Nỗi Đau Của Ngành Bán Lẻ Siêu Thị Mini",
        "content": "- Thanh toán ùn tắc: Quầy thu ngân thủ công dễ sai sót.\n- Tồn kho mất kiểm soát: Không nắm bắt được số lượng thực tế.\n- Rào cản thanh toán QR: Khách phải tự nhập số tiền, dễ sai sót.\n- Thiếu chăm sóc thành viên: Không tự động hóa tích điểm."
    },
    {
        "title": "MedMart POS - Giải Pháp Chuyển Đổi Số Toàn Diện",
        "content": "- Storefront cho Khách hàng: Đặt mua trực tuyến, thanh toán QR động tự động điền tiền.\n- Phân hệ Cashier POS: Giao diện thu ngân siêu tốc, in biên lai monospace.\n- Quản lý Kho & Sản phẩm: Thêm mới trực quan (Visual Editor), cập nhật tồn kho tức thì.\n- Thống kê & Dashboard: Biểu đồ trực quan giúp ra quyết định chính xác."
    },
    {
        "title": "Phân Tích Hệ Thống Qua Biểu Đồ Use Case",
        "content": "- Khách hàng: Xem sản phẩm, quản lý giỏ hàng online, điền SĐT tích điểm, đặt hàng & thanh toán QR.\n- Thu ngân: Đăng nhập, lập hóa đơn bán hàng tại quầy, in biên lai.\n- Quản lý: Theo dõi biểu đồ doanh thu, quản lý danh mục sản phẩm, lịch sử kho."
    },
    {
        "title": "Quy Trình Mua Hàng & Thanh Toán QR Động",
        "content": "1. Chọn hàng: Duyệt sản phẩm -> Thêm vào giỏ.\n2. Áp mã VIP: Nhập số điện thoại -> Hệ thống tự hiển thị chiết khấu.\n3. Quét VietQR: Sinh mã QR động hiển thị lập tức chính xác số tiền.\n4. Duyệt đơn: Dữ liệu đồng bộ sang màn hình hóa đơn POS nội bộ."
    },
    {
        "title": "Mô Hình Cơ Sở Dữ Liệu Vật Lý (ERD)",
        "content": "- Sản phẩm (Products): Tên, danh mục, giá bán, tồn kho, link ảnh.\n- Khách hàng (Customers): Hạng VIP/Vàng/Bạc, điểm tích lũy.\n- Hóa đơn (Invoices): Mảng sản phẩm mua, tổng tiền, phương thức thanh toán.\n- Nhật ký kho (Stock_Entries): Lưu trữ biến động kho (ai làm, số lượng)."
    },
    {
        "title": "Quản Trị Dự Án Theo Mô Hình Agile/Scrum",
        "content": "- Epic 1: eCommerce Storefront (Đặt hàng online, thanh toán QR).\n- Epic 2: Cashier POS Register (Bán hàng quầy, in biên lai).\n- Epic 3: Inventory & Products Admin (Quản lý tồn kho, chỉnh sửa ảnh).\nQuy trình: Phân rã User Stories và Tasks trong từng Sprint."
    },
    {
        "title": "Công Nghệ & Quy Trình Hợp Tác Git Workflow",
        "content": "- Tech Stack: HTML5, CSS3, Vanilla JS Module, localStorage.\n- Git Workflow:\n  + Clone dự án về local.\n  + Tạo nhánh tính năng (feature branch).\n  + Lập trình -> Commit -> Push.\n  + Tạo Pull Request để review code và gộp vào Develop."
    },
    {
        "title": "Demo Giao Diện Thực Tế Hệ Thống",
        "content": "[Chèn hình Cửa hàng Storefront online]\n\n[Chèn hình Dashboard, Biểu đồ doanh thu]\n\n[Chèn hình Form sản phẩm với chức năng Preview ảnh]"
    },
    {
        "title": "Tổng Kết & Định Hướng Phát Triển",
        "content": "Kết quả đạt được:\n- Hoàn thành 100% yêu cầu chức năng.\n- Giao diện thẩm mỹ cao (Dark mode glassmorphism).\n- Tuân thủ quy trình Git, Agile/Scrum.\n\nĐịnh hướng tương lai:\n- Tích hợp Backend hoàn chỉnh (Laravel/NodeJS)."
    }
]

for slide_data in slides_data:
    slide_layout = prs.slide_layouts[1] # 1 is title and content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = slide_data["title"]
    content.text = slide_data["content"]

prs.save("presentation_minimart.pptx")
print("Successfully generated presentation_minimart.pptx")
